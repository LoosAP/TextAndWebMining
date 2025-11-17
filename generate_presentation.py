from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from datetime import date

# Magyar, tömörebb, érdekesebb diasor
TITLE = "Text és Webbányászat – Interaktív Dashboard"
SUBTITLE = "Összefoglaló és tanulságok"

SECTIONS = [
    {
        "title": "Célok",
        "bullets": [
            "Probléma: RSS cikkek gyors feltárása és összehasonlítása",
            "Megoldás: Panel alapú, kattintható dashboard (live grafikonok)",
            "ML fókusz: K‑Means, Döntési fa, SVM, KNN összevetése",
            "Metrikák: futási idő + teljesítmény (Silhouette / Accuracy)",
        ],
    },
    {
        "title": "Adat és folyamat",
        "bullets": [
            "Forrás: articles.csv (fetch_data.py) – cím + leírás",
            "Szöveg-előkészítés: TF‑IDF (max 1000 jellemző), magyar stopwords",
            "Vizualizáció: hvPlot + Panel (kártyák, beviteli mezők)",
            "Stabilitás: random_state=42 a reprodukálhatóságért",
        ],
    },
    {
        "title": "Dashboard highlightok",
        "bullets": [
            "Kulcsszó gyakoriság források szerint (oszlopdiagram)",
            "Top szavak forrásonként (TF‑IDF, magyar stopwords)",
            "ML eredmények (idő + pontszám) – két vonaldiagram egymás mellett",
            "Extra: Andris érkezési idők 08:00–09:00 között (idősor)",
        ],
    },
    {
        "title": "ML módszertan, röviden",
        "bullets": [
            "K‑Means (k=5) → Silhouette: klaszterek szeparáltsága",
            "Pszeudo‑címkék: a K‑Means címkéit tanulja vissza Tree/SVM/KNN",
            "Train/Test: 70/30 szétválasztás TF‑IDF vektorokon",
            "Értelmezés: Accuracy nem valós címke, hanem reprodukálhatóság",
        ],
    },
    {
        "title": "Tanulságok",
        "bullets": [
            "Idő: modelltől és adatmennyiségtől függően változik",
            "Silhouette: −1..1, magasabb → jobb, elválasztható klaszterek",
            "Pontosság: klaszterstruktúra visszatanulhatósága (nem ground truth)",
            "Javítási ötletek: n‑gramok, ékezet‑normalizálás, valódi címkék",
        ],
    },
]


def _style_paragraphs_bold_lead(paragraphs):
    for p in paragraphs:
        for run in p.runs:
            run.font.size = Pt(22)
        # félkövér kiemelés a ":" előtti részre, ha van
        if p.text and ":" in p.text:
            lead, rest = p.text.split(":", 1)
            p.clear()
            r1 = p.add_run()
            r1.text = lead + ":"
            r1.font.bold = True
            r1.font.size = Pt(22)
            r1.font.color.rgb = RGBColor(59, 130, 246)  # kék árnyalat
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(22)


def add_title_slide(prs: Presentation):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = TITLE
    subtitle = slide.placeholders[1]
    subtitle.text = f"{SUBTITLE}\n{date.today().isoformat()}"
    for shape in (slide.shapes.title, subtitle):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = content.paragraphs[0]
        else:
            p = content.add_paragraph()
        p.text = b
        p.level = 0
    _style_paragraphs_bold_lead(content.paragraphs)


def add_code_slide(prs: Presentation, title: str, lines: list[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()
    for i, line in enumerate(lines):
        p = content.paragraphs[0] if i == 0 else content.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(55, 65, 81)  # szürke


def build_deck() -> Presentation:
    prs = Presentation()
    add_title_slide(prs)
    for section in SECTIONS:
        add_bullet_slide(prs, section["title"], section["bullets"])
    add_code_slide(
        prs,
        "Futtatás / Demó",
        [
            "pip install -r requirements.txt",
            "panel serve analyze_data.py --show",
            "python generate_presentation.py",
        ],
    )
    add_bullet_slide(prs, "Köszönöm!", ["Kérdések? Ötletek a továbblépéshez szívesen! "]) 
    return prs


if __name__ == "__main__":
    prs = build_deck()
    output = "project_overview.pptx"
    prs.save(output)
    print(f"Mentve: {output}")
